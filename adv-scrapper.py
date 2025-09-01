import os

# --- New Feature Dependencies ---
# This updated script can optionally extract text from PDF, Excel, and PowerPoint files.
# To use this feature, you must install the required libraries first.
# Open your terminal or command prompt and run the following command:
# pip install PyPDF2 pandas openpyxl python-pptx

# Attempt to import the necessary libraries. If they are not found,
# the script will still run but will not be able to process special files.
try:
    import PyPDF2
    import pandas as pd
    from pptx import Presentation
    LIBS_INSTALLED = True
except ImportError:
    LIBS_INSTALLED = False
    print("\n--- WARNING ---")
    print("Optional libraries not found (PyPDF2, pandas, python-pptx).")
    print("The script will skip PDF, Excel, and PowerPoint files.")
    print("To enable extraction for these files, run: pip install PyPDF2 pandas openpyxl python-pptx")
    print("---------------\n")


def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file."""
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
    except Exception as e:
        return f"[Could not extract text from PDF: {e}]"

def extract_text_from_excel(file_path):
    """Extracts data from an Excel file and formats it as CSV."""
    try:
        # Use pandas to read all sheets
        xls = pd.ExcelFile(file_path)
        text = ""
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text += f"--- Sheet: {sheet_name} ---\n"
            text += df.to_csv(index=False) + "\n"
        return text
    except Exception as e:
        return f"[Could not extract data from Excel: {e}]"

def extract_text_from_ppt(file_path):
    """Extracts text from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"[Could not extract text from PowerPoint: {e}]"


def scrape_directory_content(root_dir, output_file_path):
    """
    Traverses a directory tree, copies the content of all text files,
    and stores it in a single output file, preserving the tree structure.

    Args:
        root_dir (str): The absolute or relative path to the directory to be scraped.
        output_file_path (str): The path where the output .txt file will be saved.
    """
    # Check if the provided root directory exists.
    if not os.path.isdir(root_dir):
        print(f"Error: The directory '{root_dir}' does not exist.")
        return

    # Define file extensions to handle interactively
    interactive_extensions = {'.pdf', '.xlsx', '.xls', '.pptx', '.ppt'}

    # Define a set of common binary file extensions to always ignore.
    binary_extensions = {
        # Images
        '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg', '.tiff', '.ico', '.webp',
        # Audio
        '.mp3', '.wav', '.ogg', '.flac',
        # Video
        '.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv',
        # Archives
        '.zip', '.rar', '.7z', '.tar', '.gz',
        # Executables and compiled files
        '.exe', '.dll', '.so', '.bin', '.pyc', '.class', '.o'
    }

    try:
        # Open the output file in write mode with UTF-8 encoding.
        with open(output_file_path, 'w', encoding='utf-8') as output_file:
            # Add a header to the output file.
            output_file.write(f"Content scraped from directory: {os.path.abspath(root_dir)}\n")
            output_file.write("=" * 80 + "\n\n")

            # os.walk() generates the file names in a directory tree
            for dirpath, _, filenames in os.walk(root_dir):
                if filenames:
                    relative_dir_path = os.path.relpath(dirpath, root_dir)
                    if relative_dir_path == '.':
                        relative_dir_path = 'Root Directory'
                    output_file.write(f"// Path: {relative_dir_path}\n")
                    output_file.write("-" * 80 + "\n")

                    for filename in filenames:
                        file_path = os.path.join(dirpath, filename)
                        _, file_extension = os.path.splitext(filename)
                        file_ext_lower = file_extension.lower()

                        # --- LOGIC FOR DIFFERENT FILE TYPES ---

                        # 1. Always-skip binary files
                        if file_ext_lower in binary_extensions:
                            output_file.write(f"\n--- Binary File (content omitted): {filename} ---\n\n")
                            continue

                        # 2. Interactive binary files (PDF, Excel, etc.)
                        if LIBS_INSTALLED and file_ext_lower in interactive_extensions:
                            while True:
                                choice = input(f"Found '{filename}'. Scrape its content? (y/n): ").lower()
                                if choice in ['y', 'yes']:
                                    output_file.write(f"\n--- File (scraped content): {filename} ---\n\n")
                                    content = ""
                                    if file_ext_lower == '.pdf':
                                        content = extract_text_from_pdf(file_path)
                                    elif file_ext_lower in ['.xls', '.xlsx']:
                                        content = extract_text_from_excel(file_path)
                                    elif file_ext_lower in ['.ppt', '.pptx']:
                                        content = extract_text_from_ppt(file_path)
                                    output_file.write(content + "\n\n")
                                    break
                                elif choice in ['n', 'no']:
                                    output_file.write(f"\n--- File (skipped by user): {filename} ---\n\n")
                                    break
                                else:
                                    print("Invalid input. Please enter 'y' or 'n'.")
                            continue
                        
                        # If libraries for special files are not installed, treat them as regular binary files.
                        elif not LIBS_INSTALLED and file_ext_lower in interactive_extensions:
                             output_file.write(f"\n--- Special File (skipped, libraries not installed): {filename} ---\n\n")
                             continue

                        # 3. Regular text files
                        try:
                            with open(file_path, 'r', encoding='utf-8', errors='ignore') as current_file:
                                output_file.write(f"\n--- File: {filename} ---\n\n")
                                content = current_file.read()
                                output_file.write(content)
                                output_file.write("\n\n")
                        except Exception as e:
                            print(f"Could not read file: {file_path}. Reason: {e}")
        
        print(f"\nSuccessfully scraped all text content into '{output_file_path}'")

    except IOError as e:
        print(f"Error writing to output file '{output_file_path}'. Reason: {e}")
    except Exception as e:
        print(f"An unexpected error occurred: {e}")


if __name__ == "__main__":
    # Get the directory path from the user.
    input_directory = input("Enter the path of the directory to scrape: ")

    # Define the name of the output file.
    output_filename = "directory_content_output.txt"
    
    # Call the main function to start the process.
    scrape_directory_content(input_directory, output_filename)

